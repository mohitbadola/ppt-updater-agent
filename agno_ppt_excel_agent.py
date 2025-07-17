from agno.tools import tool
from dataclasses import dataclass
import pandas as pd
from pptx import Presentation
import re
from typing import Dict, Any

@tool(name="extract_excel_data", description="Extracts key-value mapping from an Excel file.")
@dataclass
class ExtractExcelData:
    def run(self, file_path: str) -> Dict[str, str]:
        xls = pd.ExcelFile(file_path)
        data = {}
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name)
            for col in df.columns:
                for i, val in enumerate(df[col]):
                    if pd.notnull(val):
                        key = f"{sheet_name}{col}{i}"
                        data[key] = str(val)
        return data

@tool(name="extract_ppt_text", description="Extracts text from PowerPoint slides as editable runs.")
@dataclass
class ExtractPPTText:
    def run(self, file_path: str) -> Dict[str, Any]:
        prs = Presentation(file_path)
        run_map = {}
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        for r in p.runs:
                            key = f"slide_{i}run{r.text.strip()}"
                            run_map[key] = {"slide": i, "text": r.text.strip()}
        return run_map

@tool(name="update_ppt_with_excel", description="Updates values in the PPT using new values from Excel.")
@dataclass
class UpdatePPTWithExcel:
    def run(self, pptx_path: str, excel_data: Dict[str, str], output_path: str) -> str:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        for r in p.runs:
                            current_text = r.text.strip()
                            for val in excel_data.values():
                                if (
                                    re.fullmatch(r"[\d,]+\.?\d*", current_text.replace(",", ""))
                                    and current_text != val
                                ):
                                    try:
                                        float(current_text.replace(",", ""))
                                        float(val.replace(",", ""))
                                        r.text = val
                                    except:
                                        pass
                                elif current_text == val:
                                    continue
        prs.save(output_path)
        return f"Updated PPT saved at {output_path}"